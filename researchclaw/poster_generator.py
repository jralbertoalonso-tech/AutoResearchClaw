"""Generador de Póster A0 para congresos médicos.

Crea un único slide A0 vertical (84.1 cm × 118.9 cm) con 3 columnas profesionales,
fuentes grandes para lectura a distancia y soporte para gráfico de barras automático
a partir de los datos numéricos extraídos de la sección Resultados.

Estructura del póster:
  ┌─────────────────────────────────────────────────────────┐
  │  [LOGO HOSP]      TÍTULO (72pt)        [LOGO UNIV]      │  Header ~13 cm
  │                   Autores · Institución                  │
  ├───────────────┬───────────────┬─────────────────────────┤
  │  Introducción │   Resultados  │   Discusión              │
  │  Métodos      │   + Gráfico   │   Conclusiones           │
  │               │               │   Referencias            │
  └───────────────┴───────────────┴─────────────────────────┘
      Col 1            Col 2              Col 3

Uso:
    from researchclaw.poster_generator import generate_poster
    poster_path = generate_poster(
        paper_md="...",
        output_path=Path("deliverables/poster_congreso.pptx"),
        title="Mi Investigación",
        authors="García A, López B, Martínez C",
        institution="Hospital Universitario — Oncología",
        congress="ESMO 2025",
        logos_dir=Path("assets/logos"),
    )
"""

from __future__ import annotations

import io
import re
from pathlib import Path


# ---------------------------------------------------------------------------
# Paleta de colores (misma que pptx_generator para coherencia visual)
# ---------------------------------------------------------------------------

_DARK_BLUE  = (0x0A, 0x29, 0x4B)   # cabecera y barras de sección
_MID_BLUE   = (0x1E, 0x6F, 0xA8)   # acento columnas 1 y 2
_ACCENT     = (0xC0, 0x39, 0x2B)   # acento resultados (rojo científico)
_LIGHT_BLUE = (0xD6, 0xEA, 0xF8)   # fondo de columnas
_WHITE      = (0xFF, 0xFF, 0xFF)
_DARK_GREY  = (0x1A, 0x1A, 0x2E)   # texto principal
_LIGHT_BG   = (0xF4, 0xF6, 0xF8)   # fondo general del póster
_GOLD       = (0xF3, 0x9C, 0x12)   # conclusiones take-home

# ---------------------------------------------------------------------------
# Dimensiones del póster (cm) — A0 vertical
# ---------------------------------------------------------------------------

_A0_W: float = 84.1
_A0_H: float = 118.9

_MARGIN:      float = 2.0    # márgenes laterales izquierdo/derecho
_HEADER_H:    float = 13.5   # altura de la cabecera (título + autores)
_FOOTER_H:    float = 0.5    # margen inferior de seguridad
_COL_GAP:     float = 1.2    # espacio entre columnas
_SEC_GAP:     float = 0.6    # espacio entre secciones dentro de una columna
_SEC_TITLE_H: float = 2.2    # altura de barra de título de sección
_LOGO_W:      float = 10.0   # ancho máximo de logos
_LOGO_H:      float = 8.0    # alto máximo de logos

# Columnas: 3 iguales que llenan el ancho disponible
_COL_W:  float = (_A0_W - 2 * _MARGIN - 2 * _COL_GAP) / 3   # ~25.9 cm
_COL1_X: float = _MARGIN
_COL2_X: float = _MARGIN + _COL_W + _COL_GAP
_COL3_X: float = _MARGIN + 2 * (_COL_W + _COL_GAP)

_CONTENT_Y:  float = _HEADER_H + 0.6
_CONTENT_BOT: float = _A0_H - _FOOTER_H - 0.3
_CONTENT_H:   float = _CONTENT_BOT - _CONTENT_Y   # ~104 cm disponibles en columnas

# ---------------------------------------------------------------------------
# Tamaños de fuente (pt)
# ---------------------------------------------------------------------------

_PT_TITLE   = 68
_PT_AUTHORS = 30
_PT_AFFIL   = 24
_PT_CONGRESS = 22
_PT_SEC_HEAD = 40
_PT_BODY    = 24
_PT_REF     = 18

# ---------------------------------------------------------------------------
# Palabras clave para asignación de secciones a columnas
# ---------------------------------------------------------------------------

_COL1_KEYS = frozenset([
    "antecedentes", "background", "introducción", "introduccion",
    "objetivo", "objetivos", "objectives", "purpose", "justificación",
    "justificacion", "hipótesis", "hipotesis",
    "métodos", "metodos", "methods", "material", "diseño", "diseño del estudio",
    "población", "poblacion", "pacientes", "participantes",
    "estadística", "estadistica", "análisis estadístico",
])

_COL2_KEYS = frozenset([
    "resultados", "results", "hallazgos", "findings", "datos",
    "outcomes", "desenlaces",
])

_COL3_KEYS = frozenset([
    "discusión", "discusion", "discussion",
    "conclusiones", "conclusions", "conclusión", "conclusion",
    "referencias", "references", "bibliography", "bibliografía", "bibliografia",
    "limitaciones", "limitations",
    "perspectivas", "implicaciones", "implications",
])


# ---------------------------------------------------------------------------
# Poster section mapping — consolidates full paper into 6 poster slots
# ---------------------------------------------------------------------------

# Each poster slot maps to a set of keywords that match paper section headings.
# Headings are matched case-insensitively; first match wins.
_POSTER_SLOTS: list[tuple[str, frozenset[str], int]] = [
    # (slot_title, keywords, target_column)
    ("Background", frozenset([
        "introduction", "introducción", "introduccion", "background",
        "antecedentes", "motivation", "clinical", "scientific",
        "knowledge gap", "current state", "state of evidence",
    ]), 1),
    ("Objective", frozenset([
        "objective", "objetivos", "objetivo", "objectives", "purpose",
        "aim", "aims", "review question", "hipótesis", "hipotesis",
    ]), 1),
    ("Methods", frozenset([
        "method", "methods", "métodos", "metodos", "material",
        "search strategy", "eligibility", "prisma framework",
        "picos", "data collection", "data extraction", "study selection",
        "quality assessment", "risk-of-bias", "data analysis",
        "problem formulation", "diseño",
    ]), 1),
    ("Results", frozenset([
        "results", "resultados", "hallazgos", "findings", "outcomes",
        "prisma flow", "study characteristics", "summary table",
        "included studies", "efficacy", "safety", "interpretation",
    ]), 2),
    ("Conclusions", frozenset([
        "conclusion", "conclusions", "conclusión", "conclusiones",
        "discussion", "discusión", "discusion", "implications",
        "implicaciones", "take-home", "comparison with prior",
        "heterogeneity",
    ]), 3),
    ("References", frozenset([
        "references", "referencias", "bibliography", "bibliografía",
        "bibliografia", "limitations", "limitaciones",
    ]), 3),
]

# Maximum bullets per poster section (keeps content concise)
_MAX_BULLETS_PER_SLOT = 5
# Maximum characters per bullet (poster readability at distance)
_MAX_BULLET_CHARS = 120


def _match_heading_to_slot(heading: str) -> int | None:
    """Return the slot index (0-5) for a heading, or None if no match."""
    h = heading.lower().strip()
    # Strip numbering and bold markers
    h = re.sub(r"^\d+\.\s*", "", h)
    h = h.replace("**", "").strip()
    for idx, (_, keywords, _) in enumerate(_POSTER_SLOTS):
        if any(k in h for k in keywords):
            return idx
    return None


def _truncate_bullet(text: str, max_chars: int = _MAX_BULLET_CHARS,
                     strip_citations: bool = True) -> str:
    """Truncate a bullet to max_chars, cutting at word boundary."""
    text = re.sub(r"\*{1,2}([^*]+)\*{1,2}", r"\1", text)  # strip markdown bold
    if strip_citations:
        text = re.sub(r"\[[\d,\s]+\]", "", text)  # strip citation brackets
    text = text.strip()
    if len(text) <= max_chars:
        return text
    cut = text[:max_chars].rsplit(" ", 1)[0]
    return cut.rstrip(".,;:") + "…"


def _extract_sentences(text: str, max_count: int = 5) -> list[str]:
    """Extract key sentences from prose text (no bullet markers)."""
    # Remove markdown formatting
    clean = re.sub(r"\*{1,2}([^*]+)\*{1,2}", r"\1", text)
    clean = re.sub(r"\[[\d,\s]+\]", "", clean)
    clean = re.sub(r"^#{1,3}\s+.*$", "", clean, flags=re.MULTILINE)
    # Split into sentences
    sentences = re.split(r"(?<=[.!?])\s+", clean)
    # Filter: meaningful length, no empty
    good = [s.strip() for s in sentences if 30 < len(s.strip()) < 300]
    return good[:max_count]


def _extract_results_findings(text: str, max_findings: int = 5) -> list[str]:
    """Extract concrete findings from Results text.

    Prioritizes sentences that contain quantitative data (numbers, p-values,
    percentages, comparisons) over meta-descriptions of tables and figures.

    This avoids the generic "The table includes columns: Author/Year, Design..."
    bullets that _extract_bullets produces from Markdown lists in Results.
    """
    # Clean up markdown
    clean = re.sub(r"\*{1,2}([^*]+)\*{1,2}", r"\1", text)
    clean = re.sub(r"^#{1,3}\s+.*$", "", clean, flags=re.MULTILINE)

    # Split into sentences
    sentences = re.split(r"(?<=[.!?])\s+", clean)
    sentences = [s.strip() for s in sentences if len(s.strip()) > 30]

    # Score each sentence: higher = more likely a concrete finding
    scored: list[tuple[float, str]] = []
    for s in sentences:
        score = 0.0
        s_lower = s.lower()

        # --- Positive signals: quantitative findings ---
        # Numbers with units or context (30%, p<0.05, n=20, 4500 records)
        if re.search(r"\d+\.?\d*\s*%", s):
            score += 3.0
        if re.search(r"p\s*[<>=]\s*0\.\d+", s, re.I):
            score += 3.0
        if re.search(r"\b(n\s*=\s*\d+|CI\s*\d|OR\s*[=:]|HR\s*[=:]|RR\s*[=:])", s, re.I):
            score += 2.5
        # Comparative language
        if re.search(r"\b(significant|reduction|improvement|increase|decrease|higher|lower|fewer|more)\b", s_lower):
            score += 2.0
        if re.search(r"\b(compared|versus|vs\.?|relative|than)\b", s_lower):
            score += 1.5
        # Study counts
        if re.search(r"\b\d+\s*(studi|trial|RCT|record|patient|participant|article)", s_lower):
            score += 2.0
        # Effect language
        if re.search(r"\b(efficacy|safety|adverse|benefit|risk|outcome|symptom|score)", s_lower):
            score += 1.0
        # Any number at all
        if re.search(r"\b\d+\b", s):
            score += 0.5

        # --- Negative signals: meta-descriptions ---
        # "The table includes/provides/summarizes..."
        if re.search(r"\b(table\s+(includes|provides|summarizes|shows|presents|contains))\b", s_lower):
            score -= 5.0
        # "The following columns..."
        if re.search(r"\b(following\s+columns|column\s+names?|columns?\s*:)\b", s_lower):
            score -= 5.0
        # "The diagram shows/provides..."
        if re.search(r"\b(diagram\s+(shows|provides|represents|illustrates))\b", s_lower):
            score -= 3.0
        # "This (visual|table|figure) tool/is..."
        if re.search(r"\b(visual\s+tool|visual\s+representation|figure\s+\d+\s+(shows|presents))\b", s_lower):
            score -= 3.0
        # Bullet-style table column descriptors (e.g. "- Author/Year: The name...")
        if s.strip().startswith(("-", "•")) and re.search(r":\s*(The|A)\s+(name|study|population|intervention|comparator|outcome|result|risk)", s, re.I):
            score -= 5.0
        # Lines that describe table/figure structure rather than findings
        if re.search(r"^-\s*(Author|Design|Population|Intervention|Comparator|Results?|Risk\s+of\s+Bias|Outcome|Study)\s*:", s.strip()):
            score -= 5.0
        # Process descriptions ("provides the reason", "shows the number", "includes the following")
        if re.search(r"\b(provides\s+the\s+reason|shows\s+the\s+number|includes?\s+the\s+following|also\s+provides)\b", s_lower):
            score -= 3.0
        # "This table/tool is essential/crucial/important for..."
        if re.search(r"\b(this\s+(table|tool|figure|diagram)\s+(is|was|provides|shows|summarizes|serves))\b", s_lower):
            score -= 4.0

        scored.append((score, s))

    # Sort by score descending, take top findings
    scored.sort(key=lambda x: x[0], reverse=True)

    # Only take sentences with positive score
    findings = [s for score, s in scored if score > 0][:max_findings]

    # Strip inline citations for poster readability
    findings = [re.sub(r"\s*\([^)]*(?:et al\.|20\d{2})[^)]*\)", "", f).strip()
                for f in findings]
    findings = [re.sub(r"\s*\[[\d,\s]+\]", "", f).strip() for f in findings]

    return findings


def prepare_poster_sections(
    paper_md: str,
) -> tuple[list[tuple[str, list[str], int]], str]:
    """Consolidate a full paper into 6 poster sections.

    Returns
    -------
    sections: list of (title, bullets, column_number) for the poster
    title:    extracted paper title (or empty string)
    """
    raw_sections = _parse_sections(paper_md)

    # Accumulate content per slot
    slot_texts: list[list[str]] = [[] for _ in _POSTER_SLOTS]
    paper_title = ""

    for heading, body in raw_sections:
        h_lower = heading.lower().strip()
        # Capture title
        if h_lower in ("title", "título", "titulo") and not paper_title:
            paper_title = body.split("\n")[0].strip()
            continue
        # Skip abstract (title already captures key info)
        if h_lower in ("abstract", "resumen"):
            continue

        slot_idx = _match_heading_to_slot(heading)
        if slot_idx is not None:
            slot_texts[slot_idx].append(body)

    # Build poster sections
    result: list[tuple[str, list[str], int]] = []
    for idx, (slot_title, _, col_num) in enumerate(_POSTER_SLOTS):
        combined = "\n\n".join(slot_texts[idx])
        if not combined.strip():
            continue

        # References slot: extract full citation lines instead of bullets
        if slot_title == "References":
            refs = _extract_reference_bullets(combined)
            if refs:
                result.append((slot_title, refs, col_num))
            continue

        # Results slot: use finding-specific extractor to avoid table descriptions
        if slot_title == "Results":
            bullets = _extract_results_findings(combined, max_findings=_MAX_BULLETS_PER_SLOT)
            if not bullets:
                # Fallback to generic extraction
                bullets = _extract_sentences(combined, max_count=_MAX_BULLETS_PER_SLOT)
        else:
            # Generic path: try bullets first, then sentences from prose
            bullets = _extract_bullets(combined)
            if not bullets or all(len(b) < 20 for b in bullets):
                bullets = _extract_sentences(combined, max_count=_MAX_BULLETS_PER_SLOT)

        # Truncate and limit
        bullets = [_truncate_bullet(b) for b in bullets[:_MAX_BULLETS_PER_SLOT]]
        bullets = [b for b in bullets if b]  # remove empties

        if bullets:
            result.append((slot_title, bullets, col_num))

    return result, paper_title


def _extract_reference_bullets(text: str, max_refs: int = 4) -> list[str]:
    """Extract reference entries from combined reference text.

    Handles numbered references [1], [2] etc. and bullet-style references.
    """
    refs: list[str] = []

    # Pattern 1: Numbered references like [1] Author. Title. Journal...
    numbered = re.findall(r"\[(\d+)\]\s*(.+?)(?=\[\d+\]|\Z)", text, re.DOTALL)
    if numbered:
        for num, body in numbered[:max_refs]:
            ref = body.strip().replace("\n", " ")
            ref = re.sub(r"\s+", " ", ref)
            refs.append(_truncate_bullet(f"[{num}] {ref}", max_chars=140,
                                        strip_citations=False))
        return refs

    # Pattern 2: Lines starting with a number or bullet
    for line in text.splitlines():
        line = line.strip()
        if not line or len(line) < 20:
            continue
        # Skip headings
        if line.startswith("#"):
            continue
        # Take lines that look like references (contain year, author pattern)
        if re.search(r"\b(19|20)\d{2}\b", line) and len(line) > 30:
            refs.append(_truncate_bullet(line, max_chars=140))
            if len(refs) >= max_refs:
                break

    # Fallback: just take first meaningful lines
    if not refs:
        for line in text.splitlines():
            line = line.strip()
            if len(line) > 30 and not line.startswith("#"):
                refs.append(_truncate_bullet(line, max_chars=140))
                if len(refs) >= max_refs:
                    break

    return refs


# ---------------------------------------------------------------------------
# Utilidades de parsing Markdown
# ---------------------------------------------------------------------------

def _parse_sections(text: str) -> list[tuple[str, str]]:
    """Extrae secciones (título, cuerpo) del Markdown.

    Reconoce encabezados ## y ###.
    """
    sections: list[tuple[str, str]] = []
    pattern = re.compile(r"^#{2,3}\s+(.+)$", re.MULTILINE)
    matches = list(pattern.finditer(text))
    for i, m in enumerate(matches):
        heading = m.group(1).strip()
        start   = m.end()
        end     = matches[i + 1].start() if i + 1 < len(matches) else len(text)
        body    = text[start:end].strip()
        sections.append((heading, body))
    return sections


def _extract_bullets(text: str) -> list[str]:
    """Extrae viñetas del texto Markdown."""
    bullets: list[str] = []
    for line in text.splitlines():
        stripped = line.strip()
        if stripped.startswith(("- ", "* ", "• ")):
            b = re.sub(r"^[-*•]\s+", "", stripped).strip()
            if b:
                bullets.append(b)
    if not bullets:
        # Fallback: partir por frases
        sentences = re.split(r"[.;]\s+", text)
        bullets = [s.strip() for s in sentences if len(s.strip()) > 15]
    return bullets


def _classify_section(heading: str) -> int:
    """Devuelve 1, 2 o 3 según la columna a la que pertenece la sección.

    0 = no clasificada (usará heurística de orden de aparición).
    """
    h = heading.lower().strip()
    if any(k in h for k in _COL1_KEYS):
        return 1
    if any(k in h for k in _COL2_KEYS):
        return 2
    if any(k in h for k in _COL3_KEYS):
        return 3
    return 0


def _extract_meta(paper_md: str) -> dict[str, str]:
    """Extrae título, autores, institución, congreso del Markdown."""
    meta = {
        "title":       "",
        "authors":     "",
        "institution": "",
        "congress":    "",
    }
    for line in paper_md.splitlines():
        s = line.strip()
        if not meta["title"] and s.startswith("# "):
            meta["title"] = s[2:].strip()
        elif re.match(r"\*\*Autores?\*\*[:\s]", s, re.I):
            meta["authors"] = re.sub(r"\*\*Autores?\*\*[:\s]*", "", s, flags=re.I).strip()
        elif re.match(r"\*\*Institución\*\*[:\s]", s, re.I):
            meta["institution"] = re.sub(r"\*\*Institución\*\*[:\s]*", "", s, flags=re.I).strip()
        elif re.match(r"\*\*Congreso\*\*[:\s]", s, re.I):
            meta["congress"] = re.sub(r"\*\*Congreso\*\*[:\s]*", "", s, flags=re.I).strip()
    return meta


# ---------------------------------------------------------------------------
# Detección y generación de gráfico matplotlib
# ---------------------------------------------------------------------------

def _parse_chart_line(text: str) -> tuple[list[str], list[float]] | None:
    """Parsea una línea con marcador [CHART] en formato 'Label: X% | Label: Y%'.

    También intenta parsear 'Label: X | Label: Y' con valores numéricos.
    """
    # Quitar el marcador [CHART] y los posibles símbolos de viñeta iniciales
    line = re.sub(r"\[CHART\]\s*", "", text, flags=re.I).strip()
    line = re.sub(r"^[-*•]\s+", "", line).strip()  # limpiar viñeta inicial
    # Intentar formato "Label: value% | Label: value%"
    chunks = re.split(r"\s*\|\s*", line)
    labels: list[str] = []
    values: list[float] = []
    for chunk in chunks:
        m = re.match(r"(.+?):\s*([\d.]+)", chunk.strip())
        if m:
            labels.append(m.group(1).strip()[:25])
            values.append(float(m.group(2)))
    if len(labels) >= 2:
        return labels, values
    return None


def _extract_chart_data(results_text: str) -> tuple[list[str], list[float]] | None:
    """Detecta datos numéricos comparativos en el texto de Resultados.

    Prioridad:
    1. Línea con marcador [CHART]
    2. Porcentajes en viñetas (patrón "texto: X%")
    3. Valores numéricos en viñetas con etiqueta
    """
    # 1. Marcador explícito [CHART]
    for line in results_text.splitlines():
        if "[chart]" in line.lower():
            parsed = _parse_chart_line(line)
            if parsed:
                return parsed

    # 2. Buscar porcentajes en viñetas (p.ej. "Grupo A: 4.1% vs Grupo B: 9.6%")
    vs_pattern = re.compile(
        r"([\w\s]{3,30}):\s*([\d.]+)\s*%\s*(?:vs\.?|versus)\s*([\w\s]{3,30}):\s*([\d.]+)\s*%",
        re.I,
    )
    for line in results_text.splitlines():
        m = vs_pattern.search(line)
        if m:
            return [m.group(1).strip(), m.group(3).strip()], [float(m.group(2)), float(m.group(4))]

    # 3. Acumular hasta 5 viñetas con porcentaje
    pct_pattern = re.compile(r"^[-*•]\s*(.{5,40}):\s*([\d.]+)\s*%")
    labels, values = [], []
    for line in results_text.splitlines():
        m = pct_pattern.match(line.strip())
        if m:
            labels.append(m.group(1).strip()[:25])
            values.append(float(m.group(2)))
            if len(labels) == 5:
                break
    if len(labels) >= 2:
        return labels, values

    # 4. Intentar detectar cualquier par label+número
    num_pattern = re.compile(r"^[-*•]\s*(.{5,35}?):\s*([\d.]+)")
    for line in results_text.splitlines():
        m = num_pattern.match(line.strip())
        if m:
            try:
                val = float(m.group(2))
                labels.append(m.group(1).strip()[:25])
                values.append(val)
                if len(labels) == 5:
                    break
            except ValueError:
                pass
    if len(labels) >= 2:
        return labels, values

    return None


def _make_chart_image(
    labels: list[str],
    values: list[float],
    chart_title: str = "Resultados comparativos",
    unit_hint: str = "%",
) -> io.BytesIO | None:
    """Genera un gráfico de barras horizontal con matplotlib.

    Devuelve un BytesIO con imagen PNG, o None si matplotlib no está disponible.
    """
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        import matplotlib.patches as mpatches  # noqa: F401
    except ImportError:
        return None

    bar_colors = ["#1E6FA8", "#C0392B", "#0A294B", "#F39C12", "#27AE60"]

    fig, ax = plt.subplots(figsize=(9, max(3, len(labels) * 1.1 + 0.5)))
    y_pos = list(range(len(labels)))
    colors = [bar_colors[i % len(bar_colors)] for i in range(len(labels))]

    bars = ax.barh(y_pos, values, color=colors, edgecolor="white", linewidth=0.5, height=0.6)

    # Etiquetas en las barras
    for bar, val in zip(bars, values):
        ax.text(
            bar.get_width() + max(values) * 0.02,
            bar.get_y() + bar.get_height() / 2,
            f"{val:.1f}{unit_hint}",
            va="center", ha="left",
            fontsize=13, fontweight="bold",
            color="#1A1A2E",
        )

    ax.set_yticks(y_pos)
    ax.set_yticklabels(labels, fontsize=13)
    ax.set_xlabel(f"Valor ({unit_hint})", fontsize=12)
    ax.set_title(chart_title, fontsize=14, fontweight="bold", pad=8, color="#0A294B")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.tick_params(axis="x", labelsize=11)

    # Ampliar límite derecho para etiquetas
    ax.set_xlim(0, max(values) * 1.28)

    plt.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=180, bbox_inches="tight",
                facecolor="white", edgecolor="none")
    plt.close(fig)
    buf.seek(0)
    return buf


# ---------------------------------------------------------------------------
# Utilidades python-pptx
# ---------------------------------------------------------------------------

def _rgb(r: int, g: int, b: int):
    try:
        from pptx.dml.color import RGBColor  # type: ignore[import]
        return RGBColor(r, g, b)
    except ImportError:
        return None


def _set_font(run, size_pt: int, bold: bool = False,
              italic: bool = False, color: tuple = _DARK_GREY) -> None:
    try:
        from pptx.util import Pt  # type: ignore[import]
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        run.font.italic = italic
        c = _rgb(*color)
        if c:
            run.font.color.rgb = c
    except Exception:
        pass


def _add_rect(slide, x_cm: float, y_cm: float, w_cm: float, h_cm: float,
              fill: tuple, line: bool = False):
    """Añade un rectángulo de fondo sin borde (por defecto)."""
    try:
        from pptx.util import Cm  # type: ignore[import]
        from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: F401
        shape = slide.shapes.add_shape(1, Cm(x_cm), Cm(y_cm), Cm(w_cm), Cm(h_cm))
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(*fill)
        if not line:
            shape.line.fill.background()
        return shape
    except Exception:
        return None


def _add_textbox(slide, x_cm: float, y_cm: float, w_cm: float, h_cm: float,
                 text: str, size_pt: int, bold: bool = False,
                 italic: bool = False, color: tuple = _DARK_GREY,
                 align="left", word_wrap: bool = True) -> None:
    """Añade un cuadro de texto sencillo (una sola fuente)."""
    try:
        from pptx.util import Cm  # type: ignore[import]
        from pptx.enum.text import PP_ALIGN  # type: ignore[import]
        txb = slide.shapes.add_textbox(Cm(x_cm), Cm(y_cm), Cm(w_cm), Cm(h_cm))
        tf  = txb.text_frame
        tf.word_wrap = word_wrap
        p = tf.paragraphs[0]
        if align == "center":
            p.alignment = PP_ALIGN.CENTER  # type: ignore[attr-defined]
        elif align == "right":
            p.alignment = PP_ALIGN.RIGHT  # type: ignore[attr-defined]
        run = p.add_run()
        run.text = text
        _set_font(run, size_pt, bold=bold, italic=italic, color=color)
    except Exception:
        pass


def _add_section_block(
    slide,
    x_cm: float,
    y_cm: float,
    w_cm: float,
    title: str,
    bullets: list[str],
    header_color: tuple,
    max_height_cm: float = 30.0,
    chart_img: io.BytesIO | None = None,
    chart_height_cm: float = 12.0,
) -> float:
    """Renderiza un bloque de sección (barra de título + viñetas) y devuelve la altura usada.

    Parameters
    ----------
    slide:         diapositiva python-pptx
    x_cm, y_cm:   esquina superior izquierda del bloque
    w_cm:          ancho disponible
    title:         encabezado de la sección
    bullets:       lista de textos de viñetas
    header_color:  color RGB de la barra de encabezado
    max_height_cm: altura máxima disponible para este bloque
    chart_img:     imagen matplotlib (BytesIO) opcional, insertada tras las viñetas
    chart_height_cm: altura reservada para el gráfico

    Returns
    -------
    float: altura total consumida por este bloque (cm)
    """
    from pptx.util import Cm, Pt  # type: ignore[import]
    from pptx.enum.text import PP_ALIGN  # type: ignore[import]

    used_h = 0.0

    # — Barra de encabezado de sección ——————————————————————————
    _add_rect(slide, x_cm, y_cm + used_h, w_cm, _SEC_TITLE_H, header_color)
    _add_textbox(
        slide, x_cm + 0.3, y_cm + used_h + 0.2,
        w_cm - 0.6, _SEC_TITLE_H - 0.4,
        text=title, size_pt=_PT_SEC_HEAD, bold=True,
        color=_WHITE, word_wrap=True,
    )
    used_h += _SEC_TITLE_H + 0.3

    # — Viñetas ————————————————————————————————————————————————
    if bullets:
        # Calcular espacio disponible
        avail_for_bullets = max_height_cm - used_h
        if chart_img:
            avail_for_bullets -= (chart_height_cm + 0.4)

        # Estimación de líneas: ~0.95cm por viñeta corta, más para las largas
        # Insertar todas las viñetas en un solo text frame para fluidez
        bullets_h = min(avail_for_bullets, max(5.0, len(bullets) * 1.05 + 1.0))
        bullets_h = max(bullets_h, 3.0)

        try:
            txb = slide.shapes.add_textbox(
                Cm(x_cm), Cm(y_cm + used_h),
                Cm(w_cm), Cm(bullets_h),
            )
            tf = txb.text_frame
            tf.word_wrap = True

            for i, bullet in enumerate(bullets):
                # Limpiar marcador [CHART]
                clean = re.sub(r"\[CHART\]\s*", "", bullet, flags=re.I).strip()
                clean = re.sub(r"\*{1,2}([^*]+)\*{1,2}", r"\1", clean)  # quitar **bold** Markdown

                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_before = Pt(3)
                p.space_after  = Pt(2)
                run = p.add_run()
                run.text = f"• {clean}"
                _set_font(run, _PT_BODY, color=_DARK_GREY)
        except Exception as exc:
            print(f"[poster] Error renderizando viñetas '{title}': {exc}")

        used_h += bullets_h + 0.2

    # — Gráfico matplotlib ———————————————————————————————————————
    if chart_img:
        chart_y = y_cm + used_h
        chart_h = min(chart_height_cm, max_height_cm - used_h - 0.3)
        if chart_h >= 4.0:
            try:
                slide.shapes.add_picture(
                    chart_img,
                    Cm(x_cm), Cm(chart_y),
                    Cm(w_cm), Cm(chart_h),
                )
                used_h += chart_h + 0.3
            except Exception as exc:
                print(f"[poster] Error insertando gráfico: {exc}")

    return used_h


# ---------------------------------------------------------------------------
# Función pública principal
# ---------------------------------------------------------------------------

def generate_poster(
    paper_md: str,
    output_path: Path,
    title: str = "",
    authors: str = "",
    institution: str = "",
    congress: str = "",
    logos_dir: Path | None = None,
    logo_hospital: Path | str | None = None,
    logo_university: Path | str | None = None,
    logo_congress: Path | str | None = None,
) -> Path:
    """Genera un póster A0 en formato .pptx a partir del Markdown del paper.

    Parameters
    ----------
    paper_md:         Contenido Markdown generado por el pipeline.
    output_path:      Ruta de salida del fichero .pptx.
    title:            Título del póster (se sobreescribe si paper_md contiene # h1).
    authors:          Lista de autores (se sobreescribe si paper_md tiene **Autores:**).
    institution:      Institución (ídem).
    congress:         Nombre del congreso (ídem).
    logos_dir:        Directorio fallback con hospital_logo.png / university_logo.png.
    logo_hospital:    Ruta explícita al logo del hospital (PNG). Prioridad sobre logos_dir.
    logo_university:  Ruta explícita al logo de la universidad. Prioridad sobre logos_dir.
    logo_congress:    Ruta explícita al logo del congreso / evento. Opcional.
    """
    try:
        from pptx import Presentation  # type: ignore[import]
        from pptx.util import Cm, Pt  # type: ignore[import]
        from pptx.dml.color import RGBColor  # type: ignore[import]
        from pptx.enum.text import PP_ALIGN  # type: ignore[import]
    except ImportError as exc:
        raise ImportError(
            "python-pptx no está instalado. Ejecuta: pip install python-pptx"
        ) from exc

    # ── Extraer metadatos del Markdown ───────────────────────────────────
    meta = _extract_meta(paper_md)
    effective_title       = meta["title"]       or title       or "Póster de Investigación"
    effective_authors     = meta["authors"]     or authors     or ""
    effective_institution = meta["institution"] or institution or ""
    effective_congress    = meta["congress"]    or congress    or ""

    # ── Parsear secciones y asignar columnas ─────────────────────────────
    # Use the smart consolidator to map ~30+ paper subsections → 6 poster slots
    poster_sections, extracted_title = prepare_poster_sections(paper_md)
    if extracted_title and not effective_title:
        effective_title = extracted_title

    col1: list[tuple[str, list[str]]] = []  # (título, bullets)
    col2: list[tuple[str, list[str]]] = []
    col3: list[tuple[str, list[str]]] = []
    results_text = ""

    if poster_sections:
        # Smart path: consolidated poster sections with pre-assigned columns
        for sec_title, bullets, col_num in poster_sections:
            if col_num == 1:
                col1.append((sec_title, bullets))
            elif col_num == 2:
                col2.append((sec_title, bullets))
                # Capture results text for chart extraction
                if not results_text:
                    # Re-extract from paper for chart data (needs raw text, not bullets)
                    for h, body in _parse_sections(paper_md):
                        if _match_heading_to_slot(h) == 3:  # Results slot index
                            results_text += body + "\n"
            elif col_num == 3:
                col3.append((sec_title, bullets))
    else:
        # Fallback: legacy classification for non-standard papers
        sections = _parse_sections(paper_md)
        unclassified: list[tuple[str, list[str]]] = []
        for heading, body in sections:
            bullets = _extract_bullets(body)
            col_num = _classify_section(heading)
            if col_num == 1:
                col1.append((heading, bullets))
            elif col_num == 2:
                col2.append((heading, bullets))
                if not results_text:
                    results_text = body
            elif col_num == 3:
                col3.append((heading, bullets))
            else:
                unclassified.append((heading, bullets))
        if unclassified:
            slot = 0
            for item in unclassified:
                [col1, col2, col3][slot % 3].append(item)
                slot += 1

    # Si alguna columna está vacía, insertar placeholder
    if not col1:
        col1 = [("Introducción", ["Contenido generado por ResearchClaw pipeline."])]
    if not col2:
        col2 = [("Resultados", ["Ver documento adjunto para los resultados completos."])]
    if not col3:
        col3 = [("Conclusiones", ["Ver documento adjunto para las conclusiones."])]

    # ── Gráfico de Resultados ─────────────────────────────────────────────
    chart_img: io.BytesIO | None = None
    if results_text:
        chart_data = _extract_chart_data(results_text)
        if chart_data:
            labels, values = chart_data
            # Detectar si los valores son porcentajes
            unit = "%" if max(values) <= 100 else ""
            chart_img = _make_chart_image(labels, values,
                                           chart_title="Resultados principales",
                                           unit_hint=unit)

    # ── Crear presentación A0 ─────────────────────────────────────────────
    prs = Presentation()
    prs.slide_width  = Cm(_A0_W)
    prs.slide_height = Cm(_A0_H)

    slide_layout = prs.slide_layouts[6]   # blank
    slide = prs.slides.add_slide(slide_layout)

    # — Fondo del póster ————————————————————————————————————————————
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(*_LIGHT_BG)

    # ── CABECERA ──────────────────────────────────────────────────────────
    # Barra de cabecera (fondo azul oscuro)
    _add_rect(slide, 0, 0, _A0_W, _HEADER_H, _DARK_BLUE)

    # ── Resolver logos ───────────────────────────────────────────────────
    _logos_dir = logos_dir or (Path(__file__).resolve().parent.parent / "assets" / "logos")

    def _resolve_logo(passed: Path | str | None, default_name: str) -> Path | None:
        """Devuelve la ruta de logo resuelta: ruta explícita > logos_dir > None."""
        if passed:
            p = Path(passed)
            if p.exists() and p.stat().st_size > 0:
                return p
        candidate = _logos_dir / default_name
        return candidate if candidate.exists() else None

    hosp_logo_path  = _resolve_logo(logo_hospital,  "hospital_logo.png")
    univ_logo_path  = _resolve_logo(logo_university, "university_logo.png")
    cong_logo_path  = _resolve_logo(logo_congress,   "congress_logo.png")

    # ── Layout de cabecera adaptativo ───────────────────────────────────
    # Si hay logo de congreso en el centro, usamos layout de 3 filas:
    #   ┌[hosp]────[cong(centro)]────[univ]┐
    #   │         TÍTULO GRANDE             │
    #   │       Autores · Institución       │
    # Si no hay logo congreso, layout original (logos laterales, título encima).
    LOGO_SIDE_W: float = 10.0    # hospital / universidad
    LOGO_SIDE_H: float = 8.0
    LOGO_CONG_W: float = 9.0     # congreso (ligeramente más pequeño)
    LOGO_CONG_H: float = 6.0

    has_center_logo = cong_logo_path is not None
    logo_row_h      = (LOGO_CONG_H + 0.5) if has_center_logo else LOGO_SIDE_H

    # Posición X del área de texto (entre logos laterales)
    title_x_start = _MARGIN
    title_x_end   = _A0_W - _MARGIN

    # Logo hospital (izquierda)
    if hosp_logo_path:
        try:
            slide.shapes.add_picture(
                str(hosp_logo_path),
                Cm(_MARGIN), Cm(0.5),
                Cm(LOGO_SIDE_W), Cm(LOGO_SIDE_H),
            )
            title_x_start = _MARGIN + LOGO_SIDE_W + 0.5
        except Exception as exc:
            print(f"[poster] No se pudo insertar logo hospital: {exc}")

    # Logo universidad (derecha)
    if univ_logo_path:
        try:
            slide.shapes.add_picture(
                str(univ_logo_path),
                Cm(_A0_W - _MARGIN - LOGO_SIDE_W), Cm(0.5),
                Cm(LOGO_SIDE_W), Cm(LOGO_SIDE_H),
            )
            title_x_end = _A0_W - _MARGIN - LOGO_SIDE_W - 0.5
        except Exception as exc:
            print(f"[poster] No se pudo insertar logo universidad: {exc}")

    # Logo congreso (centro, encima del título cuando hay logos laterales)
    if cong_logo_path:
        cong_x = (_A0_W - LOGO_CONG_W) / 2
        try:
            slide.shapes.add_picture(
                str(cong_logo_path),
                Cm(cong_x), Cm(0.5),
                Cm(LOGO_CONG_W), Cm(LOGO_CONG_H),
            )
        except Exception as exc:
            print(f"[poster] No se pudo insertar logo congreso: {exc}")

    title_w = max(10.0, title_x_end - title_x_start)

    # Ajuste vertical del título según presencia de logos
    if has_center_logo:
        # Con logo de congreso en centro: título va debajo de él
        title_y  = logo_row_h + 0.3
        title_h  = 3.5
        auth_y   = title_y + title_h + 0.2
        affil_y  = auth_y + 1.5
        title_pt = _PT_TITLE - 10   # ligeramente más pequeño si hay congreso
    else:
        # Sin logo congreso: título se solapa con espacio entre logos laterales
        title_y  = 0.5
        title_h  = 6.5
        auth_y   = 7.8
        affil_y  = 10.2
        title_pt = _PT_TITLE

    # Título principal
    try:
        txb = slide.shapes.add_textbox(
            Cm(title_x_start), Cm(title_y),
            Cm(title_w), Cm(title_h),
        )
        tf = txb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER  # type: ignore[attr-defined]
        run = p.add_run()
        run.text = effective_title
        _set_font(run, title_pt, bold=True, color=_WHITE)
    except Exception as exc:
        print(f"[poster] Error en título: {exc}")

    # Autores
    if effective_authors:
        _add_textbox(
            slide, title_x_start, auth_y, title_w, 1.8,
            text=effective_authors,
            size_pt=_PT_AUTHORS, bold=False,
            color=(0xB0, 0xC4, 0xDE),
            align="center",
        )

    # Institución y Congreso (texto)
    footer_header = " · ".join(filter(None, [effective_institution, effective_congress]))
    if footer_header:
        _add_textbox(
            slide, title_x_start, affil_y, title_w, 1.8,
            text=footer_header,
            size_pt=_PT_AFFIL, italic=True,
            color=(0x90, 0xA4, 0xAE),
            align="center",
        )

    # Línea separadora bajo cabecera
    _add_rect(slide, 0, _HEADER_H - 0.3, _A0_W, 0.3, _MID_BLUE)

    # ── COLUMNAS ──────────────────────────────────────────────────────────
    col_configs = [
        (col1, _COL1_X, _MID_BLUE),
        (col2, _COL2_X, _ACCENT),
        (col3, _COL3_X, _DARK_BLUE),
    ]

    for col_idx, (sections_list, col_x, header_col) in enumerate(col_configs):
        current_y = _CONTENT_Y
        avail_h   = _CONTENT_H

        for sec_idx, (sec_title, sec_bullets) in enumerate(sections_list):
            if avail_h < 4.0:
                break  # no hay espacio

            # Determinar si hay gráfico para esta sección (solo col2, primera sección)
            section_chart = None
            chart_h_reserved = 0.0
            if col_idx == 1 and sec_idx == 0 and chart_img is not None:
                section_chart = chart_img
                chart_h_reserved = 16.0  # cm reservados para el gráfico

            used = _add_section_block(
                slide,
                x_cm=col_x,
                y_cm=current_y,
                w_cm=_COL_W,
                title=sec_title,
                bullets=sec_bullets,
                header_color=header_col,
                max_height_cm=avail_h,
                chart_img=section_chart,
                chart_height_cm=chart_h_reserved,
            )

            current_y += used + _SEC_GAP
            avail_h   -= used + _SEC_GAP

    # ── LÍNEA INFERIOR ────────────────────────────────────────────────────
    _add_rect(slide, 0, _A0_H - _FOOTER_H - 0.5, _A0_W, 0.4, _MID_BLUE)

    # ── Guardar ───────────────────────────────────────────────────────────
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"[poster] Póster A0 generado: {output_path}")
    return output_path
