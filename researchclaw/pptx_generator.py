"""Generador de presentaciones PowerPoint a partir del paper final de ResearchClaw.

Convierte el Markdown del paper en una presentación de N diapositivas adaptada
a la audiencia seleccionada.  Usa python-pptx como motor de renderizado.

Uso:
    from researchclaw.pptx_generator import generate_pptx
    pptx_path = generate_pptx(
        paper_md="...",
        output_path=Path("deliverables/presentacion.pptx"),
        n_slides=12,
        audience="Comité Científico / Congresos",
        title="Mi investigación",
        sources_note="Fuente: ResearchClaw pipeline",
    )
"""

from __future__ import annotations

import re
import textwrap
from pathlib import Path


# ---------------------------------------------------------------------------
# Paleta de colores (tema científico azul oscuro)
# ---------------------------------------------------------------------------

_DARK_BLUE   = (0x0A, 0x29, 0x4B)   # fondo de portada y secciones clave
_MID_BLUE    = (0x1E, 0x6F, 0xA8)   # acento principal
_LIGHT_BLUE  = (0xD6, 0xEA, 0xF8)   # fondo de slides de contenido
_WHITE       = (0xFF, 0xFF, 0xFF)
_DARK_GREY   = (0x2C, 0x3E, 0x50)   # texto principal
_LIGHT_GREY  = (0xEC, 0xF0, 0xF1)   # separadores


# ---------------------------------------------------------------------------
# Utilidades de parsing Markdown
# ---------------------------------------------------------------------------

def _parse_sections(text: str) -> list[tuple[str, str]]:
    """Extrae secciones (título, cuerpo) del Markdown.

    Sólo reconoce encabezados ## y ###.  Devuelve lista de (heading, content).
    """
    sections: list[tuple[str, str]] = []
    pattern = re.compile(r"^#{2,3}\s+(.+)$", re.MULTILINE)
    matches = list(pattern.finditer(text))
    for i, m in enumerate(matches):
        heading = m.group(1).strip()
        start   = m.end()
        end     = matches[i + 1].start() if i + 1 < len(matches) else len(text)
        body    = text[start:end].strip()
        sections.append((heading, body))
    return sections


def _extract_bullets(text: str, max_bullets: int = 6) -> list[str]:
    """Extrae bullets del texto o parte el texto en frases cortas."""
    bullets: list[str] = []
    # Intentar extraer líneas con - o *
    for line in text.splitlines():
        line = line.strip()
        if line.startswith(("- ", "* ", "• ")):
            bullet = line.lstrip("-* •").strip()
            if bullet:
                bullets.append(bullet)
    if not bullets:
        # Partir por frases
        sentences = re.split(r"[.;]\s+", text)
        bullets = [s.strip() for s in sentences if len(s.strip()) > 15][:max_bullets]
    return bullets[:max_bullets]


def _truncate(text: str, max_chars: int = 90) -> str:
    """Trunca texto a max_chars con ellipsis."""
    return text if len(text) <= max_chars else text[:max_chars - 1] + "…"


def _wrap(text: str, width: int = 70) -> str:
    """Envuelve texto para notas del orador."""
    return "\n".join(textwrap.wrap(text, width))


# ---------------------------------------------------------------------------
# Adaptación de contenido por audiencia
# ---------------------------------------------------------------------------

_AUDIENCE_CONFIG: dict[str, dict] = {
    "Comité Científico / Congresos": {
        "font_size_body": 18,
        "simplify": False,
        "show_stats": True,
        "footer_note": "Datos con IC 95% y valor p. Ver referencias completas en el documento adjunto.",
    },
    "Colegas Médicos / Sesión Clínica": {
        "font_size_body": 20,
        "simplify": False,
        "show_stats": True,
        "footer_note": "Fuentes primarias disponibles bajo solicitud.",
    },
    "Pacientes y Familias": {
        "font_size_body": 24,
        "simplify": True,
        "show_stats": False,
        "footer_note": "Información basada en evidencia científica. Consulta siempre a tu médico.",
    },
    "Estudiantes de Medicina": {
        "font_size_body": 20,
        "simplify": False,
        "show_stats": True,
        "footer_note": "Referencias bibliográficas en el documento fuente (ResearchClaw pipeline).",
    },
}

_DEFAULT_AUDIENCE_CFG = _AUDIENCE_CONFIG["Comité Científico / Congresos"]


def _simplify_bullet(text: str) -> str:
    """Simplificación mínima de jerga médica para audiencia de pacientes."""
    replacements = {
        r"\bneoplasia maligna\b": "cáncer",
        r"\bquimioterapia\b": "medicamentos contra el cáncer",
        r"\bmetástasis\b": "extensión del cáncer",
        r"\bsupervivencia global\b": "tiempo de vida",
        r"\bIC\s*95\s*%\b": "",
        r"\bp\s*[<=>]\s*[\d.]+": "",
        r"\bHR\b": "riesgo relativo",
        r"\bOR\b": "probabilidad",
    }
    for pattern, replacement in replacements.items():
        text = re.sub(pattern, replacement, text, flags=re.IGNORECASE)
    return text.strip(" ,;.")


# ---------------------------------------------------------------------------
# Constructor de diapositivas con python-pptx
# ---------------------------------------------------------------------------

def _rgb(r: int, g: int, b: int):  # type: ignore[return]
    try:
        from pptx.dml.color import RGBColor  # type: ignore[import]
        return RGBColor(r, g, b)
    except ImportError:
        return None


def _set_font(run, size_pt: int, bold: bool = False, color: tuple = _WHITE) -> None:
    try:
        from pptx.util import Pt  # type: ignore[import]
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        c = _rgb(*color)
        if c:
            run.font.color.rgb = c
    except Exception:
        pass


def _add_slide_number(slide, slide_num: int, total: int) -> None:
    """Añade número de diapositiva en la esquina inferior derecha."""
    try:
        from pptx.util import Inches, Pt  # type: ignore[import]
        from pptx.enum.text import PP_ALIGN  # type: ignore[import]
        nb = slide.shapes.add_textbox(Inches(9.0), Inches(7.1), Inches(0.8), Inches(0.35))
        tf = nb.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT  # type: ignore[attr-defined]
        run = p.add_run()
        run.text = f"{slide_num}/{total}"
        _set_font(run, 9, bold=False, color=(0x78, 0x90, 0x9C))
    except Exception:
        pass


def _add_title_slide(prs, title: str, subtitle: str, footer: str, slide_num: int = 1, total_slides: int = 1) -> None:
    """Diapositiva de portada con fondo azul oscuro."""
    try:
        from pptx.util import Inches, Pt  # type: ignore[import]
        from pptx.dml.color import RGBColor  # type: ignore[import]
        from pptx.enum.text import PP_ALIGN  # type: ignore[import]

        slide_layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(slide_layout)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*_DARK_BLUE)

        # Título
        txBox = slide.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(8.6), Inches(2.2))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER  # type: ignore[attr-defined]
        run = p.add_run()
        run.text = _truncate(title, 120)
        _set_font(run, 32, bold=True, color=_WHITE)

        # Subtítulo
        txBox2 = slide.shapes.add_textbox(Inches(1.0), Inches(4.5), Inches(8.0), Inches(1.2))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER  # type: ignore[attr-defined]
        run2 = p2.add_run()
        run2.text = _truncate(subtitle, 100)
        _set_font(run2, 20, bold=False, color=(0xA0, 0xC4, 0xFF))

        # Footer
        txBox3 = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9.0), Inches(0.4))
        tf3 = txBox3.text_frame
        p3 = tf3.paragraphs[0]
        p3.alignment = PP_ALIGN.CENTER  # type: ignore[attr-defined]
        run3 = p3.add_run()
        run3.text = footer
        _set_font(run3, 11, bold=False, color=(0x90, 0xA4, 0xAE))

        # Número de diapositiva
        _add_slide_number(slide, slide_num, total_slides)
    except Exception:
        pass


def _add_content_slide(
    prs,
    heading: str,
    bullets: list[str],
    sources_note: str,
    audience_cfg: dict,
    slide_num: int,
    total_slides: int,
) -> None:
    """Diapositiva de contenido con título y bullets."""
    try:
        from pptx.util import Inches, Pt  # type: ignore[import]
        from pptx.dml.color import RGBColor  # type: ignore[import]
        from pptx.enum.text import PP_ALIGN  # type: ignore[import]

        slide_layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(slide_layout)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*_LIGHT_BLUE)

        # Barra lateral izquierda (acento)
        bar = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(0), Inches(0), Inches(0.18), Inches(7.5),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(*_MID_BLUE)
        bar.line.fill.background()

        # Título de la diapositiva
        txBox = slide.shapes.add_textbox(Inches(0.35), Inches(0.2), Inches(9.3), Inches(0.9))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = _truncate(heading, 80)
        _set_font(run, 26, bold=True, color=_DARK_BLUE)

        # Bullets de contenido
        body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9.0), Inches(5.0))
        tf_body = body_box.text_frame
        tf_body.word_wrap = True
        font_size = audience_cfg.get("font_size_body", 18)

        for i, bullet in enumerate(bullets[:6]):
            text = _simplify_bullet(bullet) if audience_cfg.get("simplify") else bullet
            p_b = tf_body.paragraphs[0] if i == 0 else tf_body.add_paragraph()
            p_b.space_before = Pt(4)
            run_b = p_b.add_run()
            run_b.text = f"• {_truncate(text, 100)}"
            _set_font(run_b, font_size, bold=False, color=_DARK_GREY)

        # Pie de página con fuente (guardarraíles)
        footer_box = slide.shapes.add_textbox(Inches(0.35), Inches(6.9), Inches(9.0), Inches(0.45))
        tf_footer = footer_box.text_frame
        p_f = tf_footer.paragraphs[0]
        run_f = p_f.add_run()
        footer_text = sources_note if sources_note else audience_cfg.get("footer_note", "")
        run_f.text = _truncate(footer_text, 130)
        _set_font(run_f, 9, bold=False, color=(0x78, 0x90, 0x9C))

        # Número de diapositiva
        num_box = slide.shapes.add_textbox(Inches(9.0), Inches(6.9), Inches(0.8), Inches(0.45))
        tf_num = num_box.text_frame
        p_n = tf_num.paragraphs[0]
        p_n.alignment = PP_ALIGN.RIGHT  # type: ignore[attr-defined]
        run_n = p_n.add_run()
        run_n.text = f"{slide_num}/{total_slides}"
        _set_font(run_n, 9, bold=False, color=(0x78, 0x90, 0x9C))
    except Exception:
        pass


def _add_closing_slide(prs, conclusions: list[str], take_home: str, footer: str, slide_num: int = 1, total_slides: int = 1) -> None:
    """Diapositiva final con conclusiones sobre fondo azul oscuro."""
    try:
        from pptx.util import Inches, Pt  # type: ignore[import]
        from pptx.dml.color import RGBColor  # type: ignore[import]
        from pptx.enum.text import PP_ALIGN  # type: ignore[import]

        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*_DARK_BLUE)

        # Título
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9.0), Inches(0.9))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER  # type: ignore[attr-defined]
        run = p.add_run()
        run.text = "Conclusiones Clave"
        _set_font(run, 28, bold=True, color=_WHITE)

        # Bullets de conclusiones
        body_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(8.6), Inches(3.8))
        tf_body = body_box.text_frame
        tf_body.word_wrap = True
        for i, c in enumerate(conclusions[:4]):
            p_c = tf_body.paragraphs[0] if i == 0 else tf_body.add_paragraph()
            p_c.space_before = Pt(6)
            run_c = p_c.add_run()
            run_c.text = f"+ {_truncate(c, 100)}"
            _set_font(run_c, 18, bold=False, color=_WHITE)

        # Take-home message
        th_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9.0), Inches(1.0))
        tf_th = th_box.text_frame
        tf_th.word_wrap = True
        p_th = tf_th.paragraphs[0]
        p_th.alignment = PP_ALIGN.CENTER  # type: ignore[attr-defined]
        run_th = p_th.add_run()
        run_th.text = f"> {_truncate(take_home, 110)}"
        _set_font(run_th, 16, bold=True, color=(0xFF, 0xD7, 0x00))

        # Footer
        ft_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.6), Inches(8.2), Inches(0.4))
        tf_ft = ft_box.text_frame
        p_ft = tf_ft.paragraphs[0]
        p_ft.alignment = PP_ALIGN.CENTER  # type: ignore[attr-defined]
        run_ft = p_ft.add_run()
        run_ft.text = footer
        _set_font(run_ft, 10, bold=False, color=(0x90, 0xA4, 0xAE))

        # Número de diapositiva
        _add_slide_number(slide, slide_num, total_slides)
    except Exception:
        pass


# ---------------------------------------------------------------------------
# Función pública principal
# ---------------------------------------------------------------------------

def generate_pptx(
    paper_md: str,
    output_path: Path,
    n_slides: int = 12,
    audience: str = "Comité Científico / Congresos",
    title: str = "Investigación ResearchClaw",
    sources_note: str = "",
) -> Path:
    """Genera un fichero .pptx a partir del Markdown del paper.

    Parameters
    ----------
    paper_md:
        Contenido completo del paper en Markdown.
    output_path:
        Ruta donde guardar el fichero .pptx.
    n_slides:
        Número total de diapositivas deseado (5-30).
    audience:
        Una de las claves de _AUDIENCE_CONFIG.
    title:
        Título a mostrar en la portada.
    sources_note:
        Texto de fuente para el pie de cada diapositiva (guardarraíles).
    """
    try:
        from pptx import Presentation  # type: ignore[import]
        from pptx.util import Inches  # type: ignore[import]
    except ImportError as exc:
        raise ImportError(
            "python-pptx no está instalado. Ejecuta: pip install python-pptx"
        ) from exc

    audience_cfg = _AUDIENCE_CONFIG.get(audience, _DEFAULT_AUDIENCE_CFG)
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Parsear paper ---
    sections = _parse_sections(paper_md)
    if not sections:
        # Fallback: tratar el texto completo como una sola sección
        sections = [("Resumen de la Investigación", paper_md)]

    # --- Detectar título, conclusiones y take-home del paper ---
    paper_title = title
    # Buscar primera línea con # (h1) como título
    h1 = re.search(r"^#\s+(.+)$", paper_md, re.MULTILINE)
    if h1:
        paper_title = h1.group(1).strip()

    conclusions: list[str] = []
    take_home = "Los resultados de esta investigación abren nuevas vías para la práctica clínica."
    for heading, body in sections:
        if re.search(r"conclusi[oó]n|conclusion|take.home|recomendaci", heading, re.I):
            conclusions = _extract_bullets(body, max_bullets=4)
            if conclusions:
                take_home = conclusions[0]
            break

    if not conclusions:
        # Fallback: último párrafo del paper
        paragraphs = [p.strip() for p in paper_md.split("\n\n") if len(p.strip()) > 40]
        conclusions = paragraphs[-2:] if len(paragraphs) >= 2 else paragraphs

    # --- Seleccionar secciones para n_slides (- 2 para portada + cierre) ---
    content_slots = max(1, n_slides - 2)
    if len(sections) > content_slots:
        # Distribuir: tomar secciones con mayor densidad de texto
        sections_sorted = sorted(sections, key=lambda s: len(s[1]), reverse=True)
        selected = sections_sorted[:content_slots]
        # Reordenar según aparición original
        order = {s[0]: i for i, s in enumerate(sections)}
        selected.sort(key=lambda s: order.get(s[0], 0))
    else:
        selected = sections

    total_slides = len(selected) + 2  # portada + cierre

    # --- Construir presentación ---
    subtitle = f"Generado con ResearchClaw · Audiencia: {audience}"
    footer_global = audience_cfg.get("footer_note", "ResearchClaw — IA Médica")

    _add_title_slide(prs, paper_title, subtitle, footer_global, slide_num=1, total_slides=total_slides)

    for idx, (heading, body) in enumerate(selected, start=1):
        bullets = _extract_bullets(body, max_bullets=6)
        if not bullets:
            bullets = [_truncate(body[:300], 100)]
        _add_content_slide(
            prs, heading, bullets, sources_note, audience_cfg,
            slide_num=idx + 1, total_slides=total_slides,
        )

    _add_closing_slide(prs, conclusions, take_home, footer_global, slide_num=total_slides, total_slides=total_slides)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path
